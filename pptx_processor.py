"""
PPTX Processor Module
---------------------
Handles surgical text translation in PowerPoint files.

CRITICAL: This module follows the "NO-DESTRUCTION" rule:
- Never delete/recreate slides or shapes
- Never ungroup/regroup shapes (breaks Object IDs and animations)
- Modify text IN-PLACE at the Run level to preserve formatting
"""

import logging
from typing import Optional
from pptx import Presentation
from pptx.util import Pt
from pptx.shapes.group import GroupShape
from pptx.shapes.base import BaseShape
from pptx.table import Table, _Cell
from pptx.text.text import TextFrame

from translator import Translator

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTXProcessor:
    """
    Processes PowerPoint files for translation while preserving
    all formatting, animations, and structure.
    """
    
    def __init__(self, translator: Translator):
        """
        Initialize processor with a translator instance.
        
        Args:
            translator: Configured Translator instance
        """
        self.translator = translator
        self.stats = {
            'slides_processed': 0,
            'shapes_processed': 0,
            'text_runs_translated': 0,
            'tables_processed': 0,
            'notes_translated': 0,
            'groups_traversed': 0,
            'errors': []
        }
    
    def process_file(self, input_path: str, output_path: str) -> dict:
        """
        Process a PPTX file and translate all text content.
        
        Args:
            input_path: Path to source PPTX file
            output_path: Path to save translated PPTX file
            
        Returns:
            Dictionary with processing statistics
        """
        import traceback
        
        try:
            # Load the presentation
            logger.info(f"Loading presentation from {input_path}")
            print(f"Loading: {input_path}")
            prs = Presentation(input_path)
            num_slides = len(prs.slides)
            logger.info(f"Loaded presentation with {num_slides} slides")
            print(f"Loaded {num_slides} slides")
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    logger.info(f"Processing slide {slide_idx + 1}/{num_slides}")
                    print(f"Processing slide {slide_idx + 1}/{num_slides}")
                    self._process_slide(slide)
                    self.stats['slides_processed'] += 1
                except Exception as e:
                    error_msg = f"Error on slide {slide_idx + 1}: {e}"
                    logger.error(error_msg)
                    print(error_msg)
                    print(traceback.format_exc())
                    self.stats['errors'].append(error_msg)
                    # Continue with next slide instead of failing completely
            
            # Save the translated presentation
            logger.info(f"Saving translated presentation to {output_path}")
            print(f"Saving: {output_path}")
            prs.save(output_path)
            logger.info("Save complete")
            print("Save complete")
            
            # Add translator stats
            self.stats['translator_stats'] = self.translator.get_stats()
            
            return self.stats
            
        except Exception as e:
            logger.error(f"Error processing file: {e}")
            print(f"CRITICAL ERROR: {e}")
            print(traceback.format_exc())
            self.stats['errors'].append(str(e))
            raise
    
    def _process_slide(self, slide) -> None:
        """
        Process all shapes on a slide, including speaker notes.
        """
        # Process all shapes on the slide
        for shape in slide.shapes:
            self._process_shape(shape)
        
        # Process speaker notes
        self._process_notes(slide)
    
    def _process_shape(self, shape: BaseShape) -> None:
        """
        Recursively process a shape. This is the core algorithm that
        handles groups without breaking them.
        
        CRITICAL: For grouped shapes, we recurse into the group's
        shapes WITHOUT ungrouping. This preserves Object IDs.
        """
        try:
            # Check if this is a group shape
            if isinstance(shape, GroupShape):
                logger.debug(f"Traversing group shape with {len(shape.shapes)} children")
                self.stats['groups_traversed'] += 1
                
                # RECURSIVE CALL: Process each child shape in the group
                for child_shape in shape.shapes:
                    self._process_shape(child_shape)
                return
            
            # Process table if present
            if shape.has_table:
                self._process_table(shape.table)
                return
            
            # Process text frame if present
            if shape.has_text_frame:
                self._process_text_frame(shape.text_frame)
            
            self.stats['shapes_processed'] += 1
            
        except Exception as e:
            error_msg = f"Error processing shape: {e}"
            logger.warning(error_msg)
            self.stats['errors'].append(error_msg)
    
    def _process_text_frame(self, text_frame: TextFrame) -> None:
        """
        Process a text frame by iterating through paragraphs and runs.
        
        CRITICAL: We translate at the RUN level, not the paragraph or
        text_frame level. This preserves formatting (bold, italic, color, size).
        """
        for paragraph in text_frame.paragraphs:
            self._process_paragraph(paragraph)
    
    def _process_paragraph(self, paragraph) -> None:
        """
        Process a paragraph's runs for translation.
        
        For best translation quality, we:
        1. Collect all run texts to form complete sentences
        2. Translate the combined text
        3. Redistribute back to runs proportionally
        
        This handles cases where a sentence is split across multiple runs
        with different formatting.
        """
        runs = list(paragraph.runs)
        
        if not runs:
            return
        
        # Collect text from all runs
        original_texts = [run.text for run in runs]
        combined_text = ''.join(original_texts)
        
        # Skip if empty or whitespace only
        if not combined_text.strip():
            return
        
        # Translate the combined text
        translated_text = self.translator.translate(combined_text)
        
        if translated_text == combined_text:
            # No translation occurred, skip
            return
        
        # Redistribute translated text back to runs
        # Strategy: Proportional distribution based on original lengths
        self._redistribute_text_to_runs(runs, original_texts, translated_text)
        
        self.stats['text_runs_translated'] += len(runs)
    
    def _redistribute_text_to_runs(self, runs, original_texts, translated_text) -> None:
        """
        Redistribute translated text back to runs while preserving their existence.
        
        We use a proportional approach: each run gets a portion of the translated
        text proportional to its original length.
        """
        total_original_len = sum(len(t) for t in original_texts)
        
        if total_original_len == 0:
            # Edge case: all empty runs, put all text in first run
            if runs:
                runs[0].text = translated_text
            return
        
        # Calculate positions for splitting
        translated_len = len(translated_text)
        current_pos = 0
        
        for i, (run, orig_text) in enumerate(zip(runs, original_texts)):
            if i == len(runs) - 1:
                # Last run gets the remainder
                run.text = translated_text[current_pos:]
            else:
                # Calculate proportional length
                proportion = len(orig_text) / total_original_len
                chars_for_run = int(translated_len * proportion)
                
                # Try to break at word boundary
                end_pos = current_pos + chars_for_run
                
                # Find nearest space for cleaner breaks
                if end_pos < translated_len:
                    space_pos = translated_text.rfind(' ', current_pos, end_pos + 10)
                    if space_pos > current_pos:
                        end_pos = space_pos + 1
                
                run.text = translated_text[current_pos:end_pos]
                current_pos = end_pos
    
    def _process_table(self, table: Table) -> None:
        """
        Process all cells in a table.
        
        Tables have their own structure: Table -> Row -> Cell -> TextFrame
        """
        try:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        self._process_text_frame(cell.text_frame)
            
            self.stats['tables_processed'] += 1
            
        except Exception as e:
            error_msg = f"Error processing table: {e}"
            logger.warning(error_msg)
            self.stats['errors'].append(error_msg)
    
    def _process_notes(self, slide) -> None:
        """
        Process speaker notes for a slide.
        
        Notes are accessed via slide.notes_slide.notes_text_frame
        """
        try:
            # Check if slide has notes
            if not slide.has_notes_slide:
                return
            
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                self._process_text_frame(notes_slide.notes_text_frame)
                self.stats['notes_translated'] += 1
                
        except Exception as e:
            # Notes access can fail on some slides, log but continue
            logger.debug(f"Could not access notes: {e}")
